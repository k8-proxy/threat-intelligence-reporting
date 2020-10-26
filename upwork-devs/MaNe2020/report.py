from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

import numpy as np
import matplotlib.pyplot as plt

from lxml import etree as ET
import os
import os.path
import collections


# arr = os.listdir('xmls/.')
arr = os.listdir('bucket/.')
totalNumberOfFiles = len(arr)

sanitized = 0
repaired = 0
issue = 0
clean = 0
malicious = 0
unsupported = 0

# Empty list for adding all file types that were processed
fileTypes = []
unsupportedFileTypes = []
# Occurrences of file types
cnt = collections.Counter()

legacyFiles = ['doc', 'dot', 'xls', 'xlt', 'xlm', 'ppt', 'pot', 'pps']
modernFiles = ['docx', 'dotx', 'docm', 'dotm', 'xlsx', 'xltx', 'xlsm', 'xltm', 'pptx', 'potx', 'ppsx', 'pptm', 'potm', 'ppsm']
otherFiles = ['pdf', 'jpeg', 'jpg', 'jpe', 'png', 'gif']
supportedFiles = legacyFiles + modernFiles + otherFiles
legacy = 0
modern = 0

# Parsing XML files
for file in arr:
    try:
        tree = ET.parse('bucket/'+file)
    except OSError:
        print("An exception occurred")
    root = tree.getroot()
    type = str(root.xpath("//gw:FileType/text()", namespaces={'gw': 'http://glasswall.com/namespace'})).strip('[]').replace("'", "")
    cnt[type] += 1
    # print(type)
    if type not in fileTypes:
        fileTypes.append(type)
    if root.xpath('count(//gw:SanitisationItems[@itemCount>0])', namespaces={'gw': 'http://glasswall.com/namespace'}) > 0:
        sanitized += 1
    if root.xpath('count(//gw:RemedyItems[@itemCount>0])', namespaces={'gw': 'http://glasswall.com/namespace'}) > 0:
        repaired += 1
    if root.xpath('count(//gw:IssueItems[@itemCount>0])', namespaces={'gw': 'http://glasswall.com/namespace'}) > 0:
        issue += 1
    if root.xpath('count(//gw:ContentGroups[@groupCount="1"])', namespaces={'gw': 'http://glasswall.com/namespace'}) == 1 and issue == 0:
        malicious += 1
    if sanitized == 0 and repaired == 0 and issue == 0:
        clean += 1
    if type in legacyFiles:
        legacy += 1
    elif type in modernFiles:
        modern += 1
    elif type not in supportedFiles:
        unsupported += 1


# print("Sanitized "+str(sanitized))
# print("Issue "+str(issue))
# print("Repaired "+str(repaired))
# print("Clean "+str(clean))
# print("Malicious "+str(malicious))


sanitizedPerc = round(sanitized/totalNumberOfFiles*100, 2)

# print("Legacy "+str(legacy))
# print("Modern "+str(modern))
print(str(cnt).strip("})").replace("Counter({", ""))
# print("total "+str(totalNumberOfFiles))
# print(sanitizedPerc)


for i in range(len(fileTypes)):
    if fileTypes[i] not in supportedFiles:
        unsupportedFileTypes.append(fileTypes[i])

# print(unsupportedFileTypes)
# print("Unsupported "+str(unsupported))

# Create file types occurency bar plot
# labels, values = zip(*cnt.items())
# indexes = np.arange(len(labels))
# plt.bar(indexes, values, width=1)
# plt.xticks(indexes + 0.5, labels)
# plt.show()

# Create Report Presentation
prs = Presentation('GWtemplate2.pptx')

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)

# SLIDE 1: Overall Summary
slide1 = prs.slides.add_slide(prs.slide_layouts[1])
slide1.placeholders[13].text = str(totalNumberOfFiles)
slide1.placeholders[14].text = str(sanitizedPerc)+"%"
# when file is disallowed to enter organization
# slide1.placeholders[15].text = str(fileTypes).strip('[]').replace("'", "")
slide1.placeholders[15].text = str(cnt).strip("})").replace("Counter({", "").replace("'", "")
# placeholder - probably cannot be extracted from xml
# slide1.placeholders[17].text = '100ms'

# SLIDE 2: How it works
slide2 = prs.slides.add_slide(prs.slide_layouts[3])

# SLIDE 3: Outcomes by file types
slide3 = prs.slides.add_slide(prs.slide_layouts[4])
# Allowed
allowed = totalNumberOfFiles - issue
slide3.placeholders[12].text = str(allowed)
# Disallowed ??
slide3.placeholders[13].text = str(issue)
# Clean
slide3.placeholders[14].text = str(clean)
# Remidiated
slide3.placeholders[15].text = str(repaired)
# Sanitized
slide3.placeholders[16].text = str(sanitized)
# Held - not able to rebuild
slide3.placeholders[17].text = str(issue)


# # Placeholder for emails processed
# # SLIDE 4: Emails processed
# slide4 = prs.slides.add_slide(prs.slide_layouts[5])

# SLIDE 5: Modern and legacy types
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
slide5.placeholders[12].text = str(legacy)
slide5.placeholders[13].text = str(modern)

# SLIDE 6: Unsupported file types
slide6 = prs.slides.add_slide(prs.slide_layouts[7])
# for shape in slide6.placeholders:
#  print('%d %s' % (shape.placeholder_format.idx, shape.name))
slide6.placeholders[12].text = str(unsupportedFileTypes).strip('[]').replace("'", "")
slide6.placeholders[13].text = str(unsupported)

# # SLIDE 7: Malware detection summary
# slide7 = prs.slides.add_slide(prs.slide_layouts[8])
#
# # SLIDE 8: Malicious files details
# slide8 = prs.slides.add_slide(prs.slide_layouts[9])

prs.save('Threat_Intelligence_Report.pptx')