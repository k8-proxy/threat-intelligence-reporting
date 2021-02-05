from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

import numpy as np
import matplotlib
import matplotlib.cm
import matplotlib.pyplot as plt
import pandas as pd
import squarify

from lxml import etree as ET
import os
import os.path
import collections


arr = os.listdir('bucket/.')
totalNumberOfFiles = len(arr)

sanitized = 0
repaired = 0
issue = 0
clean = 0
malicious = 0
unsupported = 0

# Empty list for adding all file types that were processed
fileTypesAll = []
unsupportedFileTypes = []
fileSizes = []
sanitizedList = collections.Counter()
remidiatedList = collections.Counter()
issuesList = collections.Counter()

# Occurrences of file types
cnt = collections.Counter()

legacyFiles = ['doc', 'dot', 'xls', 'xlt', 'xlm', 'ppt', 'pot', 'pps']
modernFiles = ['docx', 'dotx', 'docm', 'dotm', 'xlsx', 'xltx', 'xlsm', 'xltm', 'pptx', 'potx', 'ppsx', 'pptm', 'potm', 'ppsm']
otherFiles = ['pdf', 'jpeg', 'jpg', 'jpe', 'png', 'gif', 'tiff']
supportedFiles = legacyFiles + modernFiles + otherFiles
legacy = 0
modern = 0

# Parsing XML files
for file in arr:
    try:
        tree = ET.parse('bucket/'+file)
    except OSError:
        print("An exception occurred")
    root = tree.getroot()
    type = str(root.xpath("//gw:FileType/text()", namespaces={'gw': 'http://glasswall.com/namespace'})).strip('[]').replace("'", "")
    size = str(root.xpath("//gw:DocumentSummary/gw:TotalSizeInBytes/text()", namespaces={'gw': 'http://glasswall.com/namespace'})).strip('[]').replace("'", "")
    cnt[type] += 1
    fileTypesAll.append(type)
    fileSizes.append(int(size))
    if root.xpath('count(//gw:SanitisationItems[@itemCount>0])', namespaces={'gw': 'http://glasswall.com/namespace'}) > 0:
        sanitized += 1
        sanitizedList[type] += 1
    if root.xpath('count(//gw:RemedyItems[@itemCount>0])', namespaces={'gw': 'http://glasswall.com/namespace'}) > 0:
        repaired += 1
        remidiatedList[type] += 1
    if root.xpath('count(//gw:IssueItems[@itemCount>0])', namespaces={'gw': 'http://glasswall.com/namespace'}) > 0:
        issue += 1
        issuesList[type] += 1
    if root.xpath('count(//gw:ContentGroups[@groupCount="1"])', namespaces={'gw': 'http://glasswall.com/namespace'}) == 1 and issue == 0:
        malicious += 1
    if sanitized == 0 and repaired == 0 and issue == 0:
        clean += 1
    if type in legacyFiles:
        legacy += 1
    elif type in modernFiles:
        modern += 1
    elif type not in supportedFiles:
        unsupported += 1

fileTypes = list(set(fileTypesAll))
sanitizedPerc = round(sanitized/totalNumberOfFiles*100, 2)
for i in range(len(fileTypes)):
    if fileTypes[i] not in supportedFiles:
        unsupportedFileTypes.append(fileTypes[i])


# Scatter plot for file size (in Kb)) per file type
rng = np.random.RandomState(0)
colors = rng.rand(len(fileSizes))
sizes = 1000 * rng.rand(len(fileSizes))
fileSizesKb = [x / 1024 for x in fileSizes]
plt.scatter(fileSizesKb, fileTypesAll, c=colors, s=sizes, alpha=0.3, cmap='Blues')
plt.title("File Size Distribution (Kb)", fontsize=15, fontweight="bold")
plt.savefig('img/scatter.png')
plt.close()

# Bar plot for file type distribution
df = pd.DataFrame.from_records(list(dict(cnt).items()), columns=['FileType', 'count'])
# norm = matplotlib.colors.Normalize(vmin=min(df['count']), vmax=max(df['count']))
# colors = [matplotlib.cm.Blues(norm(value)) for value in df['count']]
fig = plt.gcf()
fig.set_size_inches(14, 8)
squarify.plot(label=df['FileType'], sizes=df['count'], alpha=.4)
plt.title("Requested Files by extension", fontsize=15, fontweight="bold")
plt.axis('off')
plt.savefig('img/plot.png')
plt.close()

# "stacking" bar plot per fileType, sanitized, remidiated, issue
san = pd.DataFrame.from_records(list(dict(sanitizedList).items()), columns=['FileTypes', 'santizedNum'])
rem = pd.DataFrame.from_records(list(dict(remidiatedList).items()), columns=['FileTypes', 'remidiatedNum'])
iss = pd.DataFrame.from_records(list(dict(issuesList).items()), columns=['FileTypes', 'issuesNum'])
total = san.assign(remidiatedNumm=rem['remidiatedNum'], issuesNum=iss['issuesNum'])
total = total.set_index('FileTypes')
plt.gcf().set_size_inches(20, 15)
total.plot.bar(stacked=True)
plt.title("Issue, Remediated and Sanitization Count per File Extension", fontsize=10, fontweight="bold")
plt.savefig('img/bar.png')
plt.close()


# Create Report Presentation
prs = Presentation('GWtemplate.pptx')

title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)

# SLIDE 1: Overall Summary
slide1 = prs.slides.add_slide(prs.slide_layouts[1])
slide1.placeholders[13].text = str(totalNumberOfFiles)
slide1.placeholders[14].text = str(sanitizedPerc)+"%"
# when file is disallowed to enter organization
# slide1.placeholders[15].text = str(fileTypes).strip('[]').replace("'", "")
slide1.placeholders[15].text = str(cnt).strip("})").replace("Counter({", "").replace("'", "")
# placeholder - probably cannot be extracted from xml
# slide1.placeholders[17].text = '100ms'

# # # SLIDE 1.1: Overall Summary Graphs
slide10 = prs.slides.add_slide(prs.slide_layouts[10])
slide10.placeholders[10].insert_picture('img/bar.png')
slide10.placeholders[12].insert_picture('img/plot.png')

# # # SLIDE 1.2: Overall Summary Graphs
slide11 = prs.slides.add_slide(prs.slide_layouts[11])
# for shape in slide10.placeholders:
#  print('%d %s' % (shape.placeholder_format.idx, shape.name))
slide11.placeholders[10].insert_picture('img/scatter.png')

# SLIDE 2: How it works
slide2 = prs.slides.add_slide(prs.slide_layouts[3])

# SLIDE 3: Outcomes by file types
slide3 = prs.slides.add_slide(prs.slide_layouts[4])
# Allowed
allowed = totalNumberOfFiles - issue
slide3.placeholders[12].text = str(allowed)
# Disallowed ??
slide3.placeholders[13].text = str(issue)
# Clean
slide3.placeholders[14].text = str(clean)
# Remidiated
slide3.placeholders[15].text = str(repaired)
# Sanitized
slide3.placeholders[16].text = str(sanitized)
# Held - not able to rebuild
slide3.placeholders[17].text = str(issue)


# # Placeholder for emails processed
# # SLIDE 4: Emails processed
# slide4 = prs.slides.add_slide(prs.slide_layouts[5])

# SLIDE 5: Modern and legacy types
slide5 = prs.slides.add_slide(prs.slide_layouts[6])
slide5.placeholders[12].text = str(legacy)
slide5.placeholders[13].text = str(modern)

# SLIDE 6: Unsupported file types
slide6 = prs.slides.add_slide(prs.slide_layouts[7])
slide6.placeholders[12].text = str(unsupportedFileTypes).strip('[]').replace("'", "")
slide6.placeholders[13].text = str(unsupported)

# # SLIDE 7: Malware detection summary
# slide7 = prs.slides.add_slide(prs.slide_layouts[8])
#
# # SLIDE 8: Malicious files details
# slide8 = prs.slides.add_slide(prs.slide_layouts[9])


prs.save('Threat_Intelligence_Report.pptx')